import os
from typing import Optional
from pathlib import Path
import hashlib
import csv
import json


class DocumentExtractor:
    """Extract text from various document formats."""

    @staticmethod
    def extract_text(file_path: str) -> str:
        """
        Extract text from a document based on file extension.

        Args:
            file_path: Path to the document

        Returns:
            Extracted text content

        Raises:
            ValueError: If file type is not supported
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = path.suffix.lower()

        if file_ext == '.pdf':
            return DocumentExtractor._extract_from_pdf(file_path)
        elif file_ext == '.docx':
            return DocumentExtractor._extract_from_docx(file_path)
        elif file_ext == '.doc':
            return DocumentExtractor._extract_from_doc(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return DocumentExtractor._extract_from_pptx(file_path)
        elif file_ext in ['.txt', '.md', '.text', '.markdown']:
            return DocumentExtractor._extract_from_text(file_path)
        elif file_ext in ['.html', '.htm']:
            return DocumentExtractor._extract_from_html(file_path)
        elif file_ext == '.csv':
            return DocumentExtractor._extract_from_csv(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return DocumentExtractor._extract_from_excel(file_path)
        elif file_ext == '.json':
            return DocumentExtractor._extract_from_json(file_path)
        elif file_ext == '.rtf':
            return DocumentExtractor._extract_from_rtf(file_path)
        elif file_ext == '.odt':
            return DocumentExtractor._extract_from_odt(file_path)
        elif file_ext in ['.py', '.js', '.ts', '.tsx', '.jsx', '.java', '.cpp', '.c', '.h', '.cs', '.go', '.rb', '.php', '.swift', '.kt', '.rs', '.sql', '.sh', '.bash', '.yaml', '.yml', '.xml', '.css', '.scss', '.less']:
            return DocumentExtractor._extract_from_code(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

    @staticmethod
    def _extract_from_pdf(file_path: str) -> str:
        """
        Extract text from PDF file with multiple fallback strategies.

        Tries in order:
        1. LangChain PyPDFLoader (fast, good for most PDFs)
        2. LangChain PDFPlumberLoader (better for complex layouts)
        3. pypdf2 (fallback for edge cases)
        """
        text_content = None
        errors = []

        # Strategy 1: Try LangChain PyPDFLoader first (fast and reliable)
        try:
            from langchain_community.document_loaders import PyPDFLoader

            loader = PyPDFLoader(file_path)
            documents = loader.load()

            if documents:
                text_content = "\n\n".join([doc.page_content for doc in documents if doc.page_content.strip()])
                if text_content.strip():
                    return text_content
        except Exception as e:
            errors.append(f"PyPDFLoader: {str(e)}")

        # Strategy 2: Try LangChain PDFPlumberLoader (better for tables/complex layouts)
        try:
            from langchain_community.document_loaders import PDFPlumberLoader

            loader = PDFPlumberLoader(file_path)
            documents = loader.load()

            if documents:
                text_content = "\n\n".join([doc.page_content for doc in documents if doc.page_content.strip()])
                if text_content.strip():
                    return text_content
        except Exception as e:
            errors.append(f"PDFPlumberLoader: {str(e)}")

        # Strategy 3: Fallback to raw pdfplumber
        try:
            import pdfplumber

            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                if len(pdf.pages) == 0:
                    raise Exception("PDF contains no pages")

                for page_num, page in enumerate(pdf.pages, 1):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            text_parts.append(text)
                    except Exception as page_error:
                        print(f"Warning: Could not extract page {page_num}: {str(page_error)}")
                        continue

            if text_parts:
                text_content = "\n\n".join(text_parts)
                if text_content.strip():
                    return text_content
        except Exception as e:
            errors.append(f"pdfplumber: {str(e)}")

        # Strategy 4: Final fallback to PyPDF2 (with decryption attempt)
        try:
            from PyPDF2 import PdfReader

            text_parts = []
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)

                # Try to decrypt if encrypted
                if pdf_reader.is_encrypted:
                    try:
                        pdf_reader.decrypt('')  # Try empty password first
                    except:
                        pass  # If decryption fails, continue anyway

                if len(pdf_reader.pages) == 0:
                    raise Exception("PDF contains no pages")

                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            text_parts.append(text)
                    except Exception as page_error:
                        print(f"Warning: Could not extract page {page_num}: {str(page_error)}")
                        continue

            if text_parts:
                text_content = "\n\n".join(text_parts)
                if text_content.strip():
                    return text_content
        except Exception as e:
            errors.append(f"PyPDF2: {str(e)}")

        # If we got here, all strategies failed
        if not text_content or not text_content.strip():
            error_summary = " | ".join(errors) if errors else "Unknown error"
            raise Exception(
                f"Could not extract text from PDF using any method. "
                f"The PDF might be image-based (scanned), encrypted, or corrupted. "
                f"Errors: {error_summary}"
            )

    @staticmethod
    def _extract_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)

        except ImportError:
            raise ImportError(
                "python-docx is required for DOCX extraction. "
                "Install it with: pip install python-docx"
            )
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")

    @staticmethod
    def _extract_from_doc(file_path: str) -> str:
        """Extract text from DOC file (older Word format)."""
        try:
            # Try using textract first (if available)
            try:
                import textract
                return textract.process(file_path).decode('utf-8')
            except ImportError:
                pass
            
            # Fallback: Try using antiword via subprocess (if available)
            try:
                import subprocess
                result = subprocess.run(
                    ['antiword', file_path],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                if result.returncode == 0:
                    return result.stdout
            except (FileNotFoundError, subprocess.TimeoutExpired):
                pass
            
            # Final fallback: Try using python-docx2txt or similar
            # For now, raise an informative error
            raise ImportError(
                "DOC file extraction requires additional tools. "
                "Install textract (pip install textract) or antiword for DOC support. "
                "Alternatively, convert DOC to DOCX format."
            )
        except Exception as e:
            if isinstance(e, ImportError):
                raise
            raise Exception(f"Error extracting DOC: {str(e)}")

    @staticmethod
    def _extract_from_pptx(file_path: str) -> str:
        """Extract text from PPTX or PPT file."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_runs = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
                    # Also check for tables
                    if hasattr(shape, "table"):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    text_runs.append(cell.text.strip())
            
            return "\n\n".join(text_runs) if text_runs else ""

        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint extraction. "
                "Install it with: pip install python-pptx"
            )
        except Exception as e:
            # For .ppt (older format), python-pptx might not work
            # Provide helpful error message
            if file_path.lower().endswith('.ppt'):
                raise Exception(
                    f"Error extracting PPT: {str(e)}. "
                    "Note: .ppt (older PowerPoint format) support is limited. "
                    "Please convert to .pptx format for best results."
                )
            raise Exception(f"Error extracting PPTX: {str(e)}")

    @staticmethod
    def _extract_from_rtf(file_path: str) -> str:
        """Extract text from RTF file."""
        try:
            # Try using striprtf or pyth for RTF parsing
            try:
                from striprtf.striprtf import rtf_to_text
                with open(file_path, 'rb') as f:
                    rtf_content = f.read()
                return rtf_to_text(rtf_content.decode('utf-8', errors='ignore'))
            except ImportError:
                pass
            
            # Fallback: Basic text extraction (may include RTF markup)
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                # Simple RTF text extraction (remove RTF control words)
                import re
                # Remove RTF control words and groups
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)
                text = re.sub(r'\{[^}]*\}', '', text)
                return text.strip()
        except Exception as e:
            raise Exception(f"Error extracting RTF: {str(e)}")

    @staticmethod
    def _extract_from_odt(file_path: str) -> str:
        """Extract text from ODT (OpenDocument Text) file."""
        try:
            # ODT files are ZIP archives containing XML
            import zipfile
            from xml.etree import ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as odt:
                # Read content.xml from the ODT archive
                content = odt.read('content.xml')
                root = ET.fromstring(content)
                
                # Extract text from all text nodes
                text_parts = []
                for elem in root.iter():
                    if elem.text and elem.text.strip():
                        text_parts.append(elem.text.strip())
                
                return "\n\n".join(text_parts)
        except Exception as e:
            raise Exception(f"Error extracting ODT: {str(e)}")

    @staticmethod
    def _extract_from_text(file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")

    @staticmethod
    def _extract_from_html(file_path: str) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                text = soup.get_text()
                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text = '\n'.join(chunk for chunk in chunks if chunk)
                return text
        except ImportError:
            # Fallback to basic extraction without beautifulsoup
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                # Basic HTML tag removal
                import re
                text = re.sub(r'<[^>]+>', '', content)
                return text
        except Exception as e:
            raise Exception(f"Error extracting HTML: {str(e)}")

    @staticmethod
    def _extract_from_csv(file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
                if not rows:
                    return ""
                # Convert CSV to formatted text
                text_parts = []
                headers = rows[0] if rows else []
                text_parts.append("Headers: " + ", ".join(headers))
                text_parts.append("\nData:\n")
                for row in rows[1:]:
                    text_parts.append(" | ".join(str(cell) for cell in row))
                return "\n".join(text_parts)
        except Exception as e:
            raise Exception(f"Error extracting CSV: {str(e)}")

    @staticmethod
    def _extract_from_excel(file_path: str) -> str:
        """Extract text from Excel file."""
        try:
            import pandas as pd
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")
                text_parts.append(df.to_string(index=False))
            return "\n".join(text_parts)
        except ImportError:
            raise ImportError(
                "pandas and openpyxl are required for Excel extraction. "
                "Install with: pip install pandas openpyxl"
            )
        except Exception as e:
            raise Exception(f"Error extracting Excel: {str(e)}")

    @staticmethod
    def _extract_from_json(file_path: str) -> str:
        """Extract text from JSON file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                # Pretty print JSON for better readability
                return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            raise Exception(f"Error extracting JSON: {str(e)}")

    @staticmethod
    def _extract_from_code(file_path: str) -> str:
        """Extract text from code files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                # Add file type context
                file_ext = Path(file_path).suffix
                return f"File type: {file_ext}\n\n{content}"
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
                file_ext = Path(file_path).suffix
                return f"File type: {file_ext}\n\n{content}"
        except Exception as e:
            raise Exception(f"Error reading code file: {str(e)}")

    @staticmethod
    def get_file_hash(file_path: str) -> str:
        """
        Calculate SHA-256 hash of file content.

        Args:
            file_path: Path to file

        Returns:
            Hex digest of file hash
        """
        sha256 = hashlib.sha256()
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                sha256.update(chunk)
        return sha256.hexdigest()

    @staticmethod
    def get_file_size(file_path: str) -> int:
        """
        Get file size in bytes.

        Args:
            file_path: Path to file

        Returns:
            File size in bytes
        """
        return os.path.getsize(file_path)

    @staticmethod
    def is_supported_file(filename: str) -> bool:
        """
        Check if file type is supported.

        Args:
            filename: Name of the file

        Returns:
            True if supported, False otherwise
        """
        supported_extensions = [
            # Documents
            '.pdf', '.docx', '.doc', '.txt', '.md', '.markdown', '.text',
            '.pptx', '.ppt', '.rtf', '.odt',
            # Web
            '.html', '.htm',
            # Data
            '.csv', '.xlsx', '.xls', '.json',
            # Code files
            '.py', '.js', '.ts', '.tsx', '.jsx',
            '.java', '.cpp', '.c', '.h', '.cs',
            '.go', '.rb', '.php', '.swift', '.kt', '.rs',
            '.sql', '.sh', '.bash',
            # Config/Markup
            '.yaml', '.yml', '.xml',
            # Styles
            '.css', '.scss', '.less'
        ]
        ext = Path(filename).suffix.lower()
        return ext in supported_extensions

    @staticmethod
    def clean_extracted_text(text: str) -> str:
        """
        Clean extracted text by removing common artifacts.

        Args:
            text: Raw extracted text

        Returns:
            Cleaned text
        """
        import re

        # Remove page numbers (common patterns)
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)

        # Remove excessive whitespace
        text = re.sub(r' +', ' ', text)

        # Remove excessive newlines (but keep paragraph breaks)
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)

        # Remove common header/footer artifacts
        text = re.sub(r'\n\s*(Page \d+ of \d+)\s*\n', '\n', text, flags=re.IGNORECASE)

        # Strip leading/trailing whitespace
        text = text.strip()

        return text
